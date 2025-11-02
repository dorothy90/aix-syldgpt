# %%
from pptx import Presentation
from openpyxl import load_workbook
import os
from dotenv import load_dotenv
from pathlib import Path
import base64
from openai import OpenAI
from langchain_openai import OpenAIEmbeddings
from langchain_core.documents import Document
from pymongo import MongoClient
import numpy as np

# %%
# API KEY 정보로드
load_dotenv(override=True)
embeddings_model_name = os.getenv("EMBEDDINGS_MODEL_NAME")
vl_model_name = os.getenv("VL_MODEL_NAME")
# 클라이언트 설정
text_embedder = OpenAIEmbeddings(
    model=embeddings_model_name,
    openai_api_key=os.getenv("OPENROUTER_API_KEY"),
    openai_api_base=os.getenv("OPENROUTER_BASE_URL"),
)

vision_client = OpenAI(
    api_key=os.getenv("OPENROUTER_API_KEY"),
    base_url=os.getenv("OPENROUTER_BASE_URL"),
)


# 검색/임베딩용 공통 정규화 함수 (대/소문자 무시)
def normalize_text(text: str) -> str:
    q = str(text).replace("\r\n", "\n").replace("\r", "\n")
    q = "\n".join(line.rstrip() for line in q.split("\n")).strip()
    return q.casefold()


# MongoDB 연결
mongo_client = MongoClient("mongodb://localhost:27017/")
db = mongo_client["document_vectorstore"]
collection = db["embeddings"]


# %%
# 1. PPTX 처리 함수
def extract_pptx(file_path, output_dir="output_images"):
    """PPTX에서 텍스트와 이미지 추출"""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(file_path)
    slide_data = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_images = []

        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                filename = Path(file_path).stem
                image_filename = f"{output_dir}/{filename}_slide_{i}_img_{j}.{ext}"
                with open(image_filename, "wb") as f:
                    f.write(image_bytes)
                slide_images.append(image_filename)

        slide_data.append(
            {
                "text": "\n".join(slide_text),
                "images": slide_images,
                "page_number": i + 1,
            }
        )

    return slide_data


def describe_image(image_path):
    """이미지 설명 생성"""
    with open(image_path, "rb") as img_file:
        b64_img = base64.b64encode(img_file.read()).decode()
        response = vision_client.chat.completions.create(
            model=vl_model_name,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe this image in detail."},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"},
                        },
                    ],
                }
            ],
            max_tokens=512,
        )
    return response.choices[0].message.content


def process_pptx(file_path, output_dir="output_images"):
    """PPTX를 Document 객체로 변환"""
    slide_data = extract_pptx(file_path, output_dir)

    # 이미지 설명 생성
    for slide in slide_data:
        slide["image_descriptions"] = [describe_image(img) for img in slide["images"]]

    # Document 객체 생성
    documents = []
    for slide in slide_data:
        content_parts = []

        if slide.get("text", "").strip():
            content_parts.append(f"[슬라이드 텍스트]\n{slide['text']}")

        if slide.get("image_descriptions"):
            for img_idx, img_desc in enumerate(slide["image_descriptions"]):
                content_parts.append(f"\n[이미지 {img_idx + 1} 설명]\n{img_desc}")

        page_content = "\n".join(content_parts)

        doc = Document(
            page_content=page_content,
            metadata={
                "source": file_path,
                "doc_type": "pptx",
                "page_number": slide["page_number"],
                "slide_text": slide.get("text", ""),
                "image_count": len(slide.get("images", [])),
                "image_paths": slide.get("images", []),
            },
        )
        documents.append(doc)

    return documents


# %%
# 2. Excel 처리 함수
def process_excel(file_path):
    """Excel 파일을 Document 객체로 변환 (행별 검색 최적화)"""
    import pandas as pd
    from langchain_community.document_loaders import DataFrameLoader

    # Excel 파일의 모든 시트를 읽기
    excel_file = pd.ExcelFile(file_path)
    documents = []

    for sheet_name in excel_file.sheet_names:
        # 각 시트를 DataFrame으로 읽기
        df = pd.read_excel(file_path, sheet_name=sheet_name)

        # DataFrame이 비어있지 않은 경우만 처리
        if not df.empty:
            # 모든 컬럼을 읽기 쉽게 결합 (벡터 검색용)
            df["_combined_content"] = df.apply(
                lambda row: ", ".join(
                    [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
                ),
                axis=1,
            )

            # DataFrameLoader로 문서 생성
            loader = DataFrameLoader(df, page_content_column="_combined_content")
            sheet_docs = loader.load()

            # 각 행을 순회하면서 메타데이터 보강
            for idx, doc in enumerate(sheet_docs):
                row_data = df.iloc[idx]

                # 기본 메타데이터 추가
                doc.metadata.update(
                    {
                        "source": file_path,
                        "doc_type": "excel",
                        "sheet_name": sheet_name,
                        "row_number": idx + 2,  # Excel 행 번호 (헤더 포함)
                    }
                )

                # 원본 DataFrame의 각 컬럼을 메타데이터에 추가
                # (숫자 필터링을 위해)
                for col in df.columns:
                    if col != "_combined_content":
                        value = row_data[col]
                        # NaN이 아닌 경우만 추가
                        if pd.notna(value):
                            # 숫자 타입은 그대로, 나머지는 문자열로
                            if isinstance(value, (int, float)):
                                doc.metadata[col] = float(value)
                            else:
                                doc.metadata[col] = str(value)

            documents.extend(sheet_docs)

    return documents


# 3. 일반 텍스트/PDF/Word 처리 함수
def process_text_document(file_path):
    """TXT, PDF, DOCX 등을 Document 객체로 변환"""
    from langchain_community.document_loaders import (
        TextLoader,
        PyPDFLoader,
        Docx2txtLoader,
    )
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    file_ext = Path(file_path).suffix.lower()

    # 파일 타입에 따른 로더 선택
    if file_ext == ".txt":
        loader = TextLoader(file_path, encoding="utf-8")
        doc_type = "text"
    elif file_ext == ".pdf":
        loader = PyPDFLoader(file_path)
        doc_type = "pdf"
    elif file_ext in [".docx", ".doc"]:
        loader = Docx2txtLoader(file_path)
        doc_type = "word"
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {file_ext}")

    # 문서 로드
    documents = loader.load()

    # 메타데이터 업데이트
    for doc in documents:
        doc.metadata["doc_type"] = doc_type

    # 텍스트 분할 (긴 문서의 경우)
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=250,
        chunk_overlap=50,
        separators=[r"\n---\n", "\n\n", "\n", " ", ""],
        is_separator_regex=True,
    )
    split_documents = text_splitter.split_documents(documents)

    return split_documents


# 4. Goodocs 처리 함수 (용어사전특화)
def process_goodocs(file_path):
    """Excel 파일을 Document 객체로 변환 (용어 사전 최적화)"""
    import pandas as pd
    from langchain_community.document_loaders import DataFrameLoader

    # Excel 파일의 모든 시트를 읽기
    excel_file = pd.ExcelFile(file_path)
    documents = []

    for sheet_name in excel_file.sheet_names:
        # 각 시트를 DataFrame으로 읽기
        df = pd.read_excel(file_path, sheet_name=sheet_name)

        # DataFrame이 비어있지 않은 경우만 처리
        if not df.empty:
            # 용어-설명 형태에 최적화된 content 생성
            # 첫 번째 컬럼(용어)과 두 번째 컬럼(설명)을 명확하게 구조화
            first_col = df.columns[0]  # 용어 컬럼
            second_col = (
                df.columns[1] if len(df.columns) > 1 else first_col
            )  # 설명 컬럼

            df["_combined_content"] = df.apply(
                lambda row: (
                    f"{row[first_col]}: {row[second_col]}"
                    if pd.notna(row[second_col])
                    else str(row[first_col])
                ),
                axis=1,
            )

            # DataFrameLoader로 문서 생성
            loader = DataFrameLoader(df, page_content_column="_combined_content")
            sheet_docs = loader.load()

            # 각 행을 순회하면서 메타데이터 보강
            for idx, doc in enumerate(sheet_docs):
                row_data = df.iloc[idx]

                # 기본 메타데이터
                doc.metadata.update(
                    {
                        "source": file_path,
                        "doc_type": "excel",
                        "sheet_name": sheet_name,
                        "row_number": idx + 2,
                        # 용어를 명시적으로 저장 (정확한 매칭용)
                        "term": (
                            str(row_data[first_col])
                            if pd.notna(row_data[first_col])
                            else ""
                        ),
                        "definition": (
                            str(row_data[second_col])
                            if pd.notna(row_data[second_col])
                            else ""
                        ),
                    }
                )

                # 나머지 컬럼도 메타데이터에 추가 (카테고리, 태그 등이 있을 수 있음)
                for col in df.columns:
                    if col not in [first_col, second_col, "_combined_content"]:
                        value = row_data[col]
                        if pd.notna(value):
                            if isinstance(value, (int, float)):
                                doc.metadata[col] = float(value)
                            else:
                                doc.metadata[col] = str(value)

            documents.extend(sheet_docs)

    return documents


# 5. 구조화된 JSON 처리
def process_structured_json(file_path):
    """구조화된 JSON을 DataFrame처럼 처리"""
    import json
    import pandas as pd
    import numpy as np
    from langchain_community.document_loaders import DataFrameLoader

    def _is_notna_safe(v):
        res = pd.notna(v)
        if isinstance(res, (list, pd.Series, np.ndarray)):
            return bool(np.any(res))
        return bool(res)

    def _to_text(v):
        if isinstance(v, (list, dict)):
            return json.dumps(v, ensure_ascii=False)
        return str(v)

    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    documents = []

    # JSON 배열을 DataFrame으로 변환
    if isinstance(data, list):
        # pandas json_normalize로 중첩 구조 평탄화
        df = pd.json_normalize(data)

        # 읽기 좋은 형태로 변환
        df["_combined_content"] = df.apply(
            lambda row: "\n".join(
                [
                    f"{col}: {_to_text(val)}"
                    for col, val in row.items()
                    if col != "_combined_content" and _is_notna_safe(val)
                ]
            ),
            axis=1,
        )

        # DataFrameLoader 사용
        loader = DataFrameLoader(df, page_content_column="_combined_content")
        documents = loader.load()

        # 메타데이터 보강
        for idx, doc in enumerate(documents):
            row_data = df.iloc[idx]

            doc.metadata.update(
                {
                    "source": file_path,
                    "doc_type": "json",
                    "item_index": idx,
                }
            )

            # 모든 컬럼을 메타데이터로 추가 (검색/필터링용)
            for col in df.columns:
                if col != "_combined_content":
                    value = row_data[col]
                    if _is_notna_safe(value):
                        # 리스트는 문자열로 변환
                        if isinstance(value, list):
                            doc.metadata[col] = value
                        elif isinstance(value, (int, float)):
                            doc.metadata[col] = float(value)
                        else:
                            doc.metadata[col] = str(value)

    return documents


# 6. 통합 문서 처리 및 MongoDB 저장 함수
def process_and_store_document(
    file_path, output_dir="output_images", move_after_process=True
):
    """
    모든 타입의 문서를 처리하고 MongoDB에 저장
    """
    import shutil

    file_ext = Path(file_path).suffix.lower()
    print(f"\n{'='*60}")
    print(f"파일 처리 중: {Path(file_path).name}")
    print(f"{'='*60}")

    try:
        # 파일 타입별 처리
        if file_ext == ".pptx":
            documents = process_pptx(file_path, output_dir)
        elif file_ext in [".xlsx", ".xls"]:
            documents = process_excel(file_path)
        elif file_ext in [".txt", ".pdf", ".docx", ".doc"]:
            documents = process_text_document(file_path)
        elif file_ext == ".json":
            documents = process_structured_json(file_path)
        elif file_ext in [".goodocs"]:
            documents = process_goodocs(file_path)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {file_ext}")

        print(f"추출된 문서 수: {len(documents)}")

        # MongoDB에 저장
        for idx, doc in enumerate(documents, 1):
            # 임베딩 생성 (대/소문자 무시 정규화 적용)
            embedding_vector = text_embedder.embed_query(
                normalize_text(doc.page_content)
            )

            # MongoDB 문서 구조
            mongo_doc = {
                "page_content": doc.page_content,
                "embedding": embedding_vector,
                "metadata": dict(doc.metadata),
            }

            # 저장
            result = collection.insert_one(mongo_doc)
            print(f"  ✓ 문서 {idx}/{len(documents)} 저장 완료")

        print(f"✅ {Path(file_path).name} 처리 완료!\n")

        # 처리 완료 후 파일 이동
        if move_after_process:
            # Complete_file 폴더 경로 설정
            source_path = Path(file_path)
            complete_folder = source_path.parent.parent / "Complete_file"

            # 폴더가 없으면 생성
            complete_folder.mkdir(parents=True, exist_ok=True)

            # 목적지 파일 경로
            destination_path = complete_folder / source_path.name

            # 같은 이름의 파일이 이미 있으면 타임스탬프 추가
            if destination_path.exists():
                from datetime import datetime

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                stem = destination_path.stem
                suffix = destination_path.suffix
                destination_path = complete_folder / f"{stem}_{timestamp}{suffix}"

            # 파일 이동
            shutil.move(str(source_path), str(destination_path))
            print(f"📦 파일 이동: {destination_path}")

        return len(documents)

    except Exception as e:
        print(f"❌ 오류 발생: {e}")
        raise


# %%


# %%
# 사용 예시
if __name__ == "__main__":
    # 컬렉션 초기화 (선택사항)
    # collection.delete_many({})

    # 여러 문서 처리
    files_to_process = [
        # "/Users/daehwankim/Documents/langgraph-tutorial-main/RAG_CHATBOT/files/example.pptx",
        "/Users/daehwankim/Documents/langgraph-tutorial-main/RAG_CHATBOT/files/people.json",
        "/Users/daehwankim/Documents/langgraph-tutorial-main/RAG_CHATBOT/files/sample-word-document.docx",
        # "/Users/daehwankim/Documents/langgraph-tutorial-main/RAG_CHATBOT/files/titanic.xlsx",
        # "/path/to/your/document.pdf",
        # "/path/to/your/spreadsheet.xlsx",
        # "/path/to/your/document.docx",
    ]

    total_docs = 0
    for file_path in files_to_process:
        if os.path.exists(file_path):
            docs_count = process_and_store_document(
                file_path, move_after_process=True  # False로 설정하면 이동 안 함
            )
            total_docs += docs_count
        else:
            print(f"⚠️  파일을 찾을 수 없습니다: {file_path}")

    print(f"\n{'='*60}")
    print(f"전체 처리 완료! 총 {total_docs}개 문서가 MongoDB에 저장되었습니다.")
    print(f"{'='*60}")

# %%
