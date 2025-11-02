from pptx import Presentation
import os
from dotenv import load_dotenv

# API KEY 정보로드
load_dotenv(override=True)
vl_model_name = os.getenv("VL_MODEL_NAME")


# %%
def extract_text_and_images(file_path, output_dir):
    prs = Presentation(file_path)
    slide_data = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_images = []
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                image_filename = f"{output_dir}/slide_{i}_img_{j}.{ext}"
                with open(image_filename, "wb") as f:
                    f.write(image_bytes)
                slide_images.append(image_filename)
        slide_data.append({"text": "\n".join(slide_text), "images": slide_images})
    return slide_data


slide_data = extract_text_and_images(
    "/Users/daehwankim/Documents/langgraph-tutorial-main/example.pptx", "output_images"
)

from langchain_openai import OpenAIEmbeddings
from langchain_openai import ChatOpenAI
from langchain_teddynote import logging
import os

# logging.langsmith("pptx-multimodal-embedding")
text_embedder = OpenAIEmbeddings(
    model="text-embedding-3-small",
    openai_api_key=os.getenv("OPENROUTER_API_KEY"),
    openai_api_base=os.getenv("OPENROUTER_BASE_URL"),
)

from openai import OpenAI
from PIL import Image
import base64

# ChatOpenAI은 LangChain용 LLM 클라이언트입니다
client = OpenAI(
    api_key=os.getenv("OPENROUTER_API_KEY"),
    base_url=os.getenv("OPENROUTER_BASE_URL"),
)


def describe_image(image_path):
    with open(image_path, "rb") as img_file:
        b64_img = base64.b64encode(img_file.read()).decode()
        response = client.chat.completions.create(
            model=vl_model_name,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe this image in detail."},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"},
                        },
                    ],
                }
            ],
            max_tokens=512,
        )
    return response.choices[0].message.content


# %%
# 각 슬라이드별 이미지 설명 생성
for slide in slide_data:
    print(slide["text"])
    print(slide["images"])
    slide["image_descriptions"] = [describe_image(img) for img in slide["images"]]

# %%
# RAG를 위한 Document 구조로 변환 및 임베딩
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from pathlib import Path


def create_documents_from_slides(slide_data):
    """
    슬라이드 데이터를 LangChain Document 객체로 변환합니다.

    각 슬라이드의 텍스트와 이미지 설명을 통합하여 하나의 문서로 생성합니다.
    """
    documents = []

    for slide_idx, slide in enumerate(slide_data):
        # 텍스트와 이미지 설명을 구조화된 형태로 통합
        content_parts = []

        # 슬라이드 텍스트 추가
        if slide.get("text", "").strip():
            content_parts.append(f"[슬라이드 텍스트]\n{slide['text']}")

        # 이미지 설명 추가
        if slide.get("image_descriptions"):
            for img_idx, img_desc in enumerate(slide["image_descriptions"]):
                content_parts.append(f"\n[이미지 {img_idx + 1} 설명]\n{img_desc}")

        # 통합된 콘텐츠 생성
        page_content = "\n".join(content_parts)

        # Document 객체 생성 (메타데이터에 슬라이드 정보 포함)
        doc = Document(
            page_content=page_content,
            metadata={
                "slide_number": slide_idx + 1,
                "slide_text": slide.get("text", ""),
                "image_count": len(slide.get("images", [])),
                "image_paths": slide.get("images", []),
                "image_descriptions": slide.get("image_descriptions", []),
                "source": "pptx_slide",
            },
        )
        documents.append(doc)

    return documents


# Document 객체로 변환
documents = create_documents_from_slides(slide_data)

# 텍스트 분할 (선택사항 - 슬라이드가 매우 긴 경우)
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    length_function=len,
)

# 슬라이드별로는 분할하지 않고, 필요시에만 분할 적용
# split_docs = text_splitter.split_documents(documents)

# FAISS 벡터 스토어 생성
vectorstore = FAISS.from_documents(documents=documents, embedding=text_embedder)

# 벡터 스토어 저장 (선택사항)
index_path = ".cache/pptx_faiss_index"
vectorstore.save_local(index_path)
print(f"벡터 스토어가 {index_path}에 저장되었습니다.")

# 검색 테스트
retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
query = "주제"  # 예시 쿼리
results = retriever.invoke(query)

print(f"\n쿼리: {query}")
print(f"검색된 문서 수: {len(results)}")
for i, doc in enumerate(results, 1):
    print(f"\n--- 문서 {i} ---")
    print(f"슬라이드 번호: {doc.metadata.get('slide_number')}")
    print(f"내용 (처음 200자): {doc.page_content[:200]}...")


# %%
from pymongo import MongoClient

# MongoDB 연결

mongo_client = MongoClient("mongodb://localhost:27017/")
db = mongo_client["document_vectorstore"]
collection = db["embeddings"]

# %%

# MongoDB에 저장
for idx, doc in enumerate(documents, 1):
    # 임베딩 생성
    embedding_vector = text_embedder.embed_query(doc.page_content)

    # MongoDB 문서 구조
    mongo_doc = {
        "page_content": doc.page_content,
        "embedding": embedding_vector,
        "metadata": dict(doc.metadata),
    }

    # 저장
    result = collection.insert_one(mongo_doc)
    print(f"  ✓ 문서 {idx}/{len(documents)} 저장 완료")

# %%
